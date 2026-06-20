#!/usr/bin/env python3
"""QA a generated PPTX: report structure, and optionally render pages to PNG.

Usage: qa_render.py <pptx_path> [--render p1,p2,...] [--soffice <path>]

Structure report (always): per-slide editable text-box count, picture count, and
speaker-note length — the quickest way to confirm an editable export kept real text
and that notes were injected.

Image render (with --render): converts the PPTX to PDF via LibreOffice headless,
then rasterizes the requested 1-based pages to /tmp/qa_<n>.png via PyMuPDF, so you
can eyeball how the *editable* deck actually looks (LibreOffice can re-wrap text).
"""
import sys, os, subprocess, tempfile

try:
    from pptx import Presentation
except ImportError:
    sys.exit("ERROR: python-pptx not installed (pip install python-pptx)")

pptx_path = sys.argv[1]
render_pages = None
soffice = None
if "--render" in sys.argv:
    render_pages = [int(x) for x in sys.argv[sys.argv.index("--render") + 1].split(",") if x]
if "--soffice" in sys.argv:
    soffice = sys.argv[sys.argv.index("--soffice") + 1]

prs = Presentation(pptx_path)
print(f"{os.path.basename(pptx_path)}: {len(prs.slides)} slides")
for i, s in enumerate(prs.slides, 1):
    txt = sum(1 for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip())
    pics = sum(1 for sh in s.shapes if sh.shape_type == 13)  # PICTURE
    note = len(s.notes_slide.notes_text_frame.text) if s.has_notes_slide else 0
    print(f"  slide {i}: text_boxes={txt} pictures={pics} note_chars={note}")

if not render_pages:
    sys.exit(0)

# Locate LibreOffice.
if not soffice:
    for c in ("/Applications/LibreOffice.app/Contents/MacOS/soffice",
              "/usr/bin/soffice", "/usr/bin/libreoffice"):
        if os.path.exists(c):
            soffice = c
            break
if not soffice:
    sys.exit("ERROR: LibreOffice (soffice) not found; cannot render. Pass --soffice <path>.")

try:
    import fitz  # PyMuPDF
except ImportError:
    sys.exit("ERROR: pymupdf not installed (pip install pymupdf); cannot render.")

with tempfile.TemporaryDirectory() as td:
    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", td, pptx_path],
                   capture_output=True, text=True)
    pdf = os.path.join(td, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    if not os.path.exists(pdf):
        sys.exit("ERROR: LibreOffice failed to produce a PDF for rendering.")
    doc = fitz.open(pdf)
    for p in render_pages:
        if 1 <= p <= doc.page_count:
            out = f"/tmp/qa_{p}.png"
            doc[p - 1].get_pixmap(matrix=fitz.Matrix(1.6, 1.6)).save(out)
            print(f"  rendered page {p} -> {out}")
