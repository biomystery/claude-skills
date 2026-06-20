#!/usr/bin/env python3
"""Re-inject Marp speaker notes into an editable PPTX.

Marp keeps non-directive HTML comments as presenter notes, but the experimental
`--pptx-editable` (LibreOffice) export DROPS them. This script reads the notes back
out of the source markdown and writes them onto the matching slides of the PPTX.

Usage: add_notes.py <src_md> <pptx_path>

Slide order in the markdown (split on `---`, after frontmatter) must match the deck.
Directive comments (e.g. `<!-- _class: lead -->`, `<!-- paginate: true -->`) are
skipped; only prose comments become notes.
"""
import re, sys, os

try:
    from pptx import Presentation
except ImportError:
    sys.exit("ERROR: python-pptx not installed (pip install python-pptx)")

src_md, pptx_path = sys.argv[1], sys.argv[2]
with open(src_md) as f:
    text = f.read()

# Strip YAML frontmatter (leading --- ... ---).
if text.startswith("---\n"):
    end = text.find("\n---", 4)
    if end != -1:
        text = text[end + 4:]

# Split into slides on a `---` line of its own.
slides_md = re.split(r"(?m)^---\s*$", text)

DIRECTIVE = re.compile(r"^\s*_?[\w-]+\s*:")  # _class:, paginate:, header:, etc.


def extract_note(slide_md: str) -> str:
    notes = []
    for c in re.findall(r"<!--(.*?)-->", slide_md, re.DOTALL):
        lines = [ln for ln in c.strip().splitlines()]
        nonempty = [ln for ln in lines if ln.strip()]
        # All lines look like Marp directives -> not a note.
        if nonempty and all(DIRECTIVE.match(ln) for ln in nonempty):
            continue
        body = c.strip()
        if body:
            notes.append(body)
    return "\n\n".join(notes)


notes = [extract_note(s) for s in slides_md]

prs = Presentation(pptx_path)
n_slides, n_notes = len(prs.slides), len(notes)
if n_slides != n_notes:
    print(f"Warning: {n_slides} slides but parsed {n_notes} markdown sections; "
          f"pairing the first {min(n_slides, n_notes)} in order.", file=sys.stderr)

added = 0
for slide, note in zip(prs.slides, notes):
    if note.strip():
        slide.notes_slide.notes_text_frame.text = note
        added += 1
prs.save(pptx_path)
print(f"Added notes to {added} slide(s) in {os.path.basename(pptx_path)}.")
